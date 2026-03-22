import anthropic
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

client = anthropic.Anthropic()

# Step 1: Use Claude to generate slide content
print("Generating slide content with Claude...")
response = client.messages.create(
    model="claude-sonnet-4-6",
    max_tokens=2048,
    messages=[
        {
            "role": "user",
            "content": (
                "Create content for a 5-slide presentation about renewable energy. "
                "Return ONLY a JSON array with 5 objects, each having: "
                "\"title\" (string), \"bullets\" (array of 3-4 short strings). "
                "No markdown, no explanation, just the JSON array."
            ),
        }
    ],
)

raw = response.content[0].text.strip()
slides_data = json.loads(raw)
print(f"Got {len(slides_data)} slides from Claude.")

# Step 2: Build the PPTX locally
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT     = RGBColor(0x16, 0x21, 0x3E)
HIGHLIGHT  = RGBColor(0x0F, 0x3D, 0x6E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x4C, 0xAF, 0x50)

blank_layout = prs.slide_layouts[6]  # completely blank

for i, slide_info in enumerate(slides_data):
    slide = prs.slides.add_slide(blank_layout)

    # Background rectangle
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # Accent bar on left
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    # Slide number
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4))
    tf = num_box.text_frame
    tf.text = f"{i + 1} / {len(slides_data)}"
    tf.paragraphs[0].runs[0].font.size = Pt(11)
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_info["title"]
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = GREEN

    # Divider line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = HIGHLIGHT
    line.line.fill.background()

    # Bullets
    body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(12.0), Inches(5.0))
    tf = body_box.text_frame
    tf.word_wrap = True
    for j, bullet in enumerate(slide_info["bullets"]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.runs[0].font.size = Pt(22)
        p.runs[0].font.color.rgb = WHITE
        p.space_before = Pt(8)

output_file = "renewable_energy.pptx"
prs.save(output_file)
print(f"Saved to {output_file}")
